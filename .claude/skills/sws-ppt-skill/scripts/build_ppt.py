#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_ppt.py — 주일 예배(Culto Dominical) PPT 생성기.

deck.json (구조화된 슬라이드 명세)을 입력받아, 멀리서도 잘 보이는
고대비(흰 배경 + 검정 글씨) 16:9 PowerPoint(.pptx)를 생성한다.

사용법:
    python3 build_ppt.py deck.json salida.pptx

설계 원칙
- 콘텐츠 판단(가사 추출, 메시지 파싱, RV1960 절 확보)은 호출자(Claude)가 하고,
  이 스크립트는 "일관된 렌더링"만 담당한다 → 매주 결과물의 편차를 최소화.
- 모든 텍스트는 박스 크기에 맞춰 자동으로 폰트 크기를 키우거나 줄인다(fit_size).
  넘치지 않으면서도 항상 최대한 크게 → 원거리 가독성 보장.
"""

import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 캔버스 규격 (16:9) ─────────────────────────────────────────────
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ── 기본 색상 (흰 배경 + 검정 글씨, 고대비) ────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)   # 순검정보다 살짝 부드러운 먹색(잔상 방지)
ACCENT = RGBColor(0x14, 0x2A, 0x5A)  # 딥 네이비 — 소제목/라벨/구분선용 (흰 위 고대비)
MUTED = RGBColor(0x55, 0x55, 0x55)   # 보조 텍스트(날짜 등)
GOLD = RGBColor(0xB8, 0x86, 0x0B)    # 요절 강조용 포인트

FONT_MAIN = "Arial"   # macOS/Windows 공통 가독 폰트. deck.meta.font 로 교체 가능.

# 선형(line-art) 일러스트 PNG 위치 (scripts/ 상위의 assets/illustrations/)
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ILLUS_DIR = os.path.join(SKILL_DIR, "assets", "illustrations")
# 스페인어/영어 별칭 → 파일명
ILLUS_ALIASES = {
    # ── 판화(engraving) 장면 삽화 (우선) ──
    "transfiguracion": "transfiguracion", "transfiguración": "transfiguracion",
    "monte": "transfiguracion", "montana": "transfiguracion", "gloria": "transfiguracion",
    "moises": "dos_figuras", "elias": "dos_figuras", "moises_elias": "dos_figuras",
    "comunion": "dos_figuras", "santos": "dos_figuras", "dos_figuras": "dos_figuras",
    "nube": "nube_voz", "voz": "nube_voz", "nube_voz": "nube_voz",
    # ── 단순 라인아트 아이콘(폴백/보조) ──
    "cruz": "cross", "luz": "sun", "sol": "sun", "mountain": "transfiguracion",
    "cloud_voice": "nube_voz", "two_figures": "dos_figuras",
    "paloma": "dove", "espiritu": "dove", "biblia": "book", "libro": "book",
    "tumba": "empty_tomb", "resurreccion": "empty_tomb", "corona": "crown",
    "rey": "crown", "fuego": "flame", "llama": "flame", "corazon": "heart",
    "amor": "heart", "estrella": "star",
}


def illus_png(motivo):
    """모티프명(별칭 허용) → 라인아트 PNG 경로. 없으면 None."""
    if not motivo:
        return None
    name = ILLUS_ALIASES.get(str(motivo).lower(), str(motivo).lower())
    path = os.path.join(ILLUS_DIR, name + ".png")
    return path if os.path.exists(path) else None


# ══════════════════════════════════════════════════════════════════
# 유틸리티
# ══════════════════════════════════════════════════════════════════

def hex_to_rgb(value):
    """'#142A5A' → RGBColor. None/빈값이면 None."""
    if not value:
        return None
    v = value.lstrip("#")
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def fit_size(text, box_w_in, box_h_in, max_pt, min_pt,
             avg_char_w=0.58, line_h=1.30):
    """
    텍스트가 박스 안에 들어가는 최대 폰트 크기(pt)를 추정해 반환.
    - avg_char_w: 폰트 대비 평균 글자 폭 비율(em). 한글/스페인어 혼용 안전값.
    - line_h: 줄 간격 배수.
    명시적 줄바꿈(\n)과 자동 줄바꿈을 함께 고려한다.
    """
    if not text:
        return max_pt
    paragraphs = str(text).split("\n")
    for pt in range(int(max_pt), int(min_pt) - 1, -1):
        chars_per_line = max(1, int((box_w_in * 72.0) / (pt * avg_char_w)))
        line_px = pt * line_h
        max_lines = max(1, int((box_h_in * 72.0) / line_px))
        used = 0
        for para in paragraphs:
            length = len(para) if para else 1
            used += max(1, -(-length // chars_per_line))  # ceil 나눗셈
        if used <= max_lines:
            return pt
    return int(min_pt)


def new_slide(prs, bg=WHITE):
    """빈 레이아웃 슬라이드 + 배경색 채우기."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 완전 빈 레이아웃
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def add_text(slide, left, top, width, height, text, size, color=BLACK,
             bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_MAIN, line_spacing=1.15, italic=False):
    """텍스트 박스 추가. text 안의 \n 은 문단으로 분리."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para_text in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = para_text
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.name = font
        f.color.rgb = color
    return box


def add_rule(slide, left, top, width, color=ACCENT, height_pt=4):
    """수평 구분선(가는 막대)."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Pt(height_pt))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


# ══════════════════════════════════════════════════════════════════
# 삽화(illustration) — 외부 이미지 없이 도형으로 그리는 간단 모티프
# ══════════════════════════════════════════════════════════════════

def _shape(slide, kind, l, t, w, h, color=ACCENT, line=None):
    sp = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp


def motif_cross(slide, cx, cy, s, color=ACCENT):
    _shape(slide, MSO_SHAPE.RECTANGLE, cx - s*0.11, cy - s*0.5, s*0.22, s, color)
    _shape(slide, MSO_SHAPE.RECTANGLE, cx - s*0.38, cy - s*0.18, s*0.76, s*0.2, color)


def motif_book(slide, cx, cy, s, color=ACCENT):
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx - s*0.5, cy - s*0.35, s*0.48, s*0.7, color)
    _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + s*0.02, cy - s*0.35, s*0.48, s*0.7, color)
    _shape(slide, MSO_SHAPE.RECTANGLE, cx - s*0.02, cy - s*0.35, s*0.04, s*0.7, WHITE)


def motif_heart(slide, cx, cy, s, color=ACCENT):
    _shape(slide, MSO_SHAPE.HEART, cx - s*0.5, cy - s*0.5, s, s, color)


def motif_light(slide, cx, cy, s, color=GOLD):
    _shape(slide, MSO_SHAPE.SUN, cx - s*0.5, cy - s*0.5, s, s, color)


def motif_star(slide, cx, cy, s, color=GOLD):
    _shape(slide, MSO_SHAPE.STAR_5_POINT, cx - s*0.5, cy - s*0.5, s, s, color)


def motif_dove(slide, cx, cy, s, color=ACCENT):
    _shape(slide, MSO_SHAPE.OVAL, cx - s*0.45, cy - s*0.18, s*0.7, s*0.32, color)
    _shape(slide, MSO_SHAPE.ISOCELES_TRIANGLE, cx + s*0.1, cy - s*0.42, s*0.42, s*0.42, color)
    _shape(slide, MSO_SHAPE.OVAL, cx - s*0.5, cy - s*0.28, s*0.16, s*0.16, color)


def motif_fish(slide, cx, cy, s, color=ACCENT):
    _shape(slide, MSO_SHAPE.OVAL, cx - s*0.5, cy - s*0.22, s*0.8, s*0.44, color)
    _shape(slide, MSO_SHAPE.ISOCELES_TRIANGLE, cx + s*0.22, cy - s*0.22, s*0.34, s*0.44, color)


def motif_badge(slide, cx, cy, s, color=ACCENT):
    _shape(slide, MSO_SHAPE.OVAL, cx - s*0.5, cy - s*0.5, s, s, color)


MOTIFS = {
    "cross": motif_cross, "cruz": motif_cross,
    "book": motif_book, "biblia": motif_book, "libro": motif_book,
    "heart": motif_heart, "corazon": motif_heart, "amor": motif_heart,
    "light": motif_light, "luz": motif_light, "sol": motif_light,
    "star": motif_star, "estrella": motif_star,
    "dove": motif_dove, "paloma": motif_dove, "espiritu": motif_dove,
    "fish": motif_fish, "pez": motif_fish, "pescador": motif_fish,
}


# ══════════════════════════════════════════════════════════════════
# 슬라이드 렌더러
# ══════════════════════════════════════════════════════════════════

def render_cover(prs, s, ctx):
    slide = new_slide(prs)
    add_text(slide, 1, 2.1, SLIDE_W_IN - 2, 2.0,
             s.get("title", "Culto Dominical"), 66, ACCENT, bold=True)
    add_rule(slide, SLIDE_W_IN/2 - 2, 4.35, 4)
    sub = s.get("subtitle", "")
    if sub:
        add_text(slide, 1, 4.6, SLIDE_W_IN - 2, 0.9, sub, 30, MUTED)
    g = s.get("guiador")
    if g:
        add_text(slide, 1, 5.7, SLIDE_W_IN - 2, 0.9,
                 "Guiador:  " + g, 30, BLACK, bold=True)


def render_section(prs, s, ctx):
    """예배 순서 구분 슬라이드: 큰 순서명 + 세부 + 섬기는 사람."""
    slide = new_slide(prs)
    orden = s.get("orden", "")
    detalle = s.get("detalle", "")
    servidor = s.get("servidor", "")
    if not (detalle or servidor):
        # 부가정보 없음: 제목만 정중앙, 구분선 없음(겹침 방지)
        size = fit_size(orden, 11.5, 3.4, 90, 40)
        add_text(slide, 0.9, 2.05, SLIDE_W_IN - 1.8, 3.4, orden, size,
                 ACCENT, bold=True)
        return
    # 부가정보 있음: 제목을 상단 영역에 하단정렬 → 구분선/부가정보와 겹치지 않음
    size = fit_size(orden, 11.5, 2.3, 80, 40)
    add_text(slide, 0.9, 1.3, SLIDE_W_IN - 1.8, 2.3, orden, size,
             ACCENT, bold=True, anchor=MSO_ANCHOR.BOTTOM)
    y = 4.05
    add_rule(slide, SLIDE_W_IN/2 - 1.6, y, 3.2)
    y += 0.42
    if detalle:
        add_text(slide, 0.9, y, SLIDE_W_IN - 1.8, 1.0, detalle, 40, BLACK, bold=True)
        y += 1.05
    if servidor:
        add_text(slide, 0.9, y, SLIDE_W_IN - 1.8, 1.0, servidor, 44, BLACK, bold=True)


def render_hymn_title(prs, s, ctx):
    slide = new_slide(prs)
    num = s.get("numero", "")
    titulo = s.get("titulo", "")
    label = ("Himno  N° " + str(num)) if num else "Himno"
    add_text(slide, 1, 2.2, SLIDE_W_IN - 2, 1.2, label, 40, MUTED, bold=True)
    size = fit_size(titulo, 11.5, 2.4, 72, 34)
    add_text(slide, 1, 3.4, SLIDE_W_IN - 2, 2.4, titulo, size, ACCENT, bold=True)


def render_hymn(prs, s, ctx):
    """찬송 가사 슬라이드. heading(라벨) + body(가사) + 선택적 coro(후렴)."""
    slide = new_slide(prs)
    heading = s.get("heading", "")
    body = s.get("body", "")
    coro = s.get("coro", "")
    top = 0.5
    if heading:
        add_text(slide, 0.8, top, SLIDE_W_IN - 1.6, 0.8, heading, 30, ACCENT,
                 bold=True, anchor=MSO_ANCHOR.TOP)
        top = 1.35
    body_h = (SLIDE_H_IN - top - 0.4) * (0.62 if coro else 1.0)
    size = fit_size(body, 11.6, body_h, 54, 26)
    add_text(slide, 0.85, top, SLIDE_W_IN - 1.7, body_h, body, size, BLACK,
             bold=True, line_spacing=1.18)
    if coro:
        cy = top + body_h + 0.1
        add_rule(slide, 0.85, cy, SLIDE_W_IN - 1.7, GOLD, height_pt=3)
        c_h = SLIDE_H_IN - cy - 0.35
        csize = fit_size(coro, 11.6, c_h - 0.1, 44, 22)
        add_text(slide, 0.85, cy + 0.12, SLIDE_W_IN - 1.7, c_h - 0.12,
                 coro, csize, ACCENT, bold=True, italic=True, line_spacing=1.14)


def render_big_text(prs, s, ctx):
    """Credo / Padre Nuestro 등 지정 텍스트를 크게 꽉 채우는 슬라이드."""
    slide = new_slide(prs)
    heading = s.get("heading", "")
    texto = s.get("texto", "")
    pagina = s.get("pagina")
    top = 0.6
    if heading:
        lab = heading + (("  ({}/{})".format(pagina, s["total"]))
                         if pagina and s.get("total") else "")
        add_text(slide, 0.8, top, SLIDE_W_IN - 1.6, 0.9, lab, 30, ACCENT,
                 bold=True, anchor=MSO_ANCHOR.TOP)
        top = 1.5
    body_h = SLIDE_H_IN - top - 0.5
    size = fit_size(texto, 11.8, body_h, 50, 24)
    add_text(slide, 0.75, top, SLIDE_W_IN - 1.5, body_h, texto, size, BLACK,
             bold=True, line_spacing=1.24)


def render_passage(prs, s, ctx):
    """성경 본문(읽기) 슬라이드: 절 번호 + 절 내용, 2~4절/페이지."""
    slide = new_slide(prs)
    ref = s.get("referencia", "")
    versiculos = s.get("versiculos", [])
    add_text(slide, 0.8, 0.45, SLIDE_W_IN - 1.6, 0.85, ref, 34, ACCENT,
             bold=True, anchor=MSO_ANCHOR.TOP)
    add_rule(slide, 0.85, 1.35, SLIDE_W_IN - 1.7)
    # 절 텍스트 조립 (절번호 위첨자 대신 굵은 접두)
    lines = []
    for v in versiculos:
        n = v.get("n", "")
        t = v.get("t", "")
        lines.append("{}  {}".format(n, t) if n != "" else t)
    body = "\n".join(lines)
    body_h = SLIDE_H_IN - 1.6 - 0.45
    n_verses = max(1, len(versiculos))
    hi = 42 if n_verses <= 2 else 34
    size = fit_size(body, 11.6, body_h, hi, 20)
    add_text(slide, 0.85, 1.6, SLIDE_W_IN - 1.7, body_h, body, size, BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)


def render_scripture(prs, s, ctx):
    """메시지 중 '본문 읽기' 절: RV1960 전문. passage 와 동일 렌더 + 라벨."""
    render_passage(prs, s, ctx)


def render_message_title(prs, s, ctx):
    slide = new_slide(prs)
    kicker = s.get("kicker", "Mensaje")
    titulo = s.get("titulo", "")
    clave_ref = s.get("versiculo_clave", "")
    clave_txt = s.get("texto_clave", "")
    predicador = s.get("predicador", "")
    if kicker:
        add_text(slide, 1, 0.85, SLIDE_W_IN - 2, 0.8, kicker, 28, MUTED, bold=True)
    size = fit_size(titulo, 11.4, 2.1, 66, 32)
    add_text(slide, 0.9, 1.55, SLIDE_W_IN - 1.8, 2.1, titulo, size, ACCENT,
             bold=True, anchor=MSO_ANCHOR.BOTTOM)
    y = 3.95
    add_rule(slide, SLIDE_W_IN/2 - 1.8, y, 3.6, GOLD)
    if clave_ref or clave_txt:
        combo = (("«{}»".format(clave_txt)) if clave_txt else "")
        if clave_ref:
            combo = (combo + "\n" + clave_ref) if combo else clave_ref
        csize = fit_size(combo, 11.2, 2.0, 34, 20)
        add_text(slide, 1, y + 0.3, SLIDE_W_IN - 2, 2.0, combo, csize, BLACK,
                 italic=True, line_spacing=1.15, anchor=MSO_ANCHOR.TOP)
    if predicador:
        add_text(slide, 1, 6.55, SLIDE_W_IN - 2, 0.7, predicador, 26, MUTED, bold=True)


def render_message_part(prs, s, ctx):
    slide = new_slide(prs)
    parte = s.get("parte", "")
    subtitulo = s.get("subtitulo", "")
    if not subtitulo:
        size = fit_size(parte, 11.5, 3.4, 84, 40)
        add_text(slide, 0.9, 2.05, SLIDE_W_IN - 1.8, 3.4, parte, size,
                 ACCENT, bold=True)
        return
    size = fit_size(parte, 11.5, 2.3, 76, 40)
    add_text(slide, 0.9, 1.4, SLIDE_W_IN - 1.8, 2.3, parte, size,
             ACCENT, bold=True, anchor=MSO_ANCHOR.BOTTOM)
    add_rule(slide, SLIDE_W_IN/2 - 1.6, 4.2, 3.2)
    add_text(slide, 0.9, 4.55, SLIDE_W_IN - 1.8, 1.6, subtitulo, 36, BLACK,
             anchor=MSO_ANCHOR.TOP)


def render_illustration(prs, s, ctx):
    slide = new_slide(prs)
    caption = s.get("caption", "")
    # 1순위: 명시적 이미지 경로, 2순위: 모티프 라인아트 PNG
    img = s.get("image") or illus_png(s.get("motivo"))
    placed = False
    if img and os.path.exists(img):
        try:
            pic = slide.shapes.add_picture(img, 0, 0, height=Inches(4.4))
            pic.left = int((prs.slide_width - pic.width) / 2)
            pic.top = Inches(0.95)
            placed = True
        except Exception:
            placed = False
    if not placed:
        # PNG 없을 때만 도형 폴백
        motivo = (s.get("motivo") or "cross").lower()
        fn = MOTIFS.get(motivo, motif_badge)
        color = hex_to_rgb(s.get("color")) or (GOLD if motivo in
                ("light", "luz", "sol", "star", "estrella") else ACCENT)
        fn(slide, SLIDE_W_IN/2, 3.2, 3.2, color)
    if caption:
        size = fit_size(caption, 11.5, 1.3, 44, 24)
        add_text(slide, 0.9, 5.75, SLIDE_W_IN - 1.8, 1.35, caption, size,
                 BLACK, bold=True)


RENDERERS = {
    "cover": render_cover,
    "section": render_section,
    "hymn_title": render_hymn_title,
    "hymn": render_hymn,
    "credo": lambda prs, s, ctx: render_big_text(
        prs, {**s, "heading": s.get("heading", "Credo Apostólico")}, ctx),
    "lords_prayer": lambda prs, s, ctx: render_big_text(
        prs, {**s, "heading": s.get("heading", "Padre Nuestro")}, ctx),
    "big_text": render_big_text,
    "passage": render_passage,
    "scripture": render_scripture,
    "message_title": render_message_title,
    "key_verse_repeat": render_message_title,
    "message_part": render_message_part,
    "illustration": render_illustration,
}


def main():
    if len(sys.argv) < 3:
        print("uso: python3 build_ppt.py deck.json salida.pptx")
        sys.exit(1)
    deck_path, out_path = sys.argv[1], sys.argv[2]
    with open(deck_path, encoding="utf-8") as fh:
        deck = json.load(fh)

    meta = deck.get("meta", {})
    global FONT_MAIN, ACCENT, BLACK
    FONT_MAIN = meta.get("font", FONT_MAIN)
    if meta.get("accent"):
        ACCENT = hex_to_rgb(meta["accent"])
    if meta.get("text"):
        BLACK = hex_to_rgb(meta["text"])

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    ctx = {"meta": meta}
    count = 0
    for s in deck.get("slides", []):
        renderer = RENDERERS.get(s.get("type"))
        if renderer is None:
            print("  ! 알 수 없는 슬라이드 타입 건너뜀:", s.get("type"))
            continue
        renderer(prs, s, ctx)
        count += 1

    prs.save(out_path)
    print("생성 완료: {}  (슬라이드 {}장)".format(out_path, count))


if __name__ == "__main__":
    main()
